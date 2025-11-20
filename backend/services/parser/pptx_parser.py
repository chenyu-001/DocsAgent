"""PPTX document parser"""
from pptx import Presentation
from utils.text_clean import clean_text


class PPTXParser:
    """Parser for PowerPoint presentations"""

    def parse(self, file_path: str):
        """Extract text and metadata from PPTX file"""
        prs = Presentation(file_path)
        text_parts = []

        # Extract text from each slide
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                text_parts.append(f"[Slide {slide_num + 1}]\n" + "\n".join(slide_text))

        full_text = "\n\n".join(text_parts)
        cleaned_text = clean_text(full_text)

        # Extract metadata
        core_props = prs.core_properties
        metadata = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
        }

        return {
            "text": cleaned_text,
            "metadata": metadata,
            "page_count": len(prs.slides),
            "word_count": len(cleaned_text.split()),
        }
